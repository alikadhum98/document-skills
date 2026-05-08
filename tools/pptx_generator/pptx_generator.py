"""
مولّد عروض PowerPoint للمفاوضات والمؤتمرات - العراق
التشغيل: python tools/pptx_generator/pptx_generator.py
"""

import sys, os
sys.path.insert(0, os.path.dirname(__file__))
from pptx_data import PRESENTATIONS, PRESENTATION_TYPE, COMMON

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import datetime

# ========== الألوان ==========
COLORS = {
    "أزرق_داكن":  RGBColor(0x1F, 0x49, 0x7D),
    "أزرق":       RGBColor(0x2E, 0x75, 0xB6),
    "أزرق_فاتح":  RGBColor(0xDE, 0xEB, 0xF7),
    "أخضر":       RGBColor(0x37, 0x86, 0x44),
    "أخضر_فاتح":  RGBColor(0xE2, 0xEF, 0xDA),
    "أحمر":       RGBColor(0xC0, 0x00, 0x00),
    "أحمر_فاتح":  RGBColor(0xFF, 0xE7, 0xE7),
    "ذهبي":       RGBColor(0xBF, 0x8F, 0x00),
    "ذهبي_فاتح":  RGBColor(0xFF, 0xF2, 0xCC),
    "رمادي":      RGBColor(0x70, 0x70, 0x70),
    "أبيض":       RGBColor(0xFF, 0xFF, 0xFF),
    "نص":         RGBColor(0x26, 0x26, 0x26),
    "خلفية":      RGBColor(0xF8, 0xF9, 0xFF),
}

BULLET_COLORS = {
    "أخضر": (COLORS["أخضر"],       COLORS["أخضر_فاتح"]),
    "أزرق": (COLORS["أزرق"],       COLORS["أزرق_فاتح"]),
    "أحمر": (COLORS["أحمر"],       COLORS["أحمر_فاتح"]),
    "ذهبي": (COLORS["ذهبي"],       COLORS["ذهبي_فاتح"]),
}

CARD_COLORS = {
    "أخضر": (COLORS["أخضر"],       COLORS["أخضر_فاتح"]),
    "أزرق": (COLORS["أزرق_داكن"],  COLORS["أزرق_فاتح"]),
    "ذهبي": (COLORS["ذهبي"],       COLORS["ذهبي_فاتح"]),
    "أحمر": (COLORS["أحمر"],       COLORS["أحمر_فاتح"]),
}

FONT = "Simplified Arabic"

# ========== دوال مساعدة ==========
def set_bg(slide, color):
    bg   = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def set_rtl_para(para):
    pPr = para._p.get_or_add_pPr()
    pPr.set('rtl', '1')
    para.alignment = PP_ALIGN.RIGHT

def add_textbox(slide, text, left, top, width, height,
                size=18, bold=False, color=None, align=PP_ALIGN.RIGHT,
                bg_color=None, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    if bg_color:
        txb.fill.solid()
        txb.fill.fore_color.rgb = bg_color
        txb.line.fill.background()
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    set_rtl_para(p)
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name  = FONT
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.color.rgb = color or COLORS["نص"]
    return txb

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width     = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape

def add_slide_number(slide, n, total):
    add_textbox(slide, f"{n}/{total}",
                Emu(300000), Emu(4800000), Emu(600000), Emu(280000),
                size=11, color=COLORS["رمادي"], align=PP_ALIGN.LEFT)

# ========== بناء الشرائح ==========
def build_title_slide(prs, pres_data, common):
    blank  = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(blank)
    set_bg(slide, COLORS["أزرق_داكن"])

    # شريط زخرفي علوي
    add_rect(slide, Emu(0), Emu(0), Emu(9144000), Emu(180000), COLORS["أزرق"])

    # مستطيل المحتوى
    add_rect(slide, Emu(300000), Emu(600000),
             Emu(8544000), Emu(3000000), RGBColor(0x17, 0x37, 0x6A))

    # الأيقونة الديكورية
    add_textbox(slide, "🌍", Emu(4000000), Emu(700000),
                Emu(1000000), Emu(800000), size=40,
                align=PP_ALIGN.CENTER)

    # العنوان
    add_textbox(slide, pres_data["العنوان"],
                Emu(400000), Emu(1600000), Emu(8344000), Emu(1000000),
                size=30, bold=True, color=COLORS["أبيض"],
                align=PP_ALIGN.CENTER)

    # العنوان الفرعي
    add_textbox(slide, pres_data["العنوان_الفرعي"],
                Emu(400000), Emu(2700000), Emu(8344000), Emu(600000),
                size=18, color=COLORS["أزرق_فاتح"],
                align=PP_ALIGN.CENTER)

    # شريط المعلومات
    add_rect(slide, Emu(0), Emu(3900000), Emu(9144000), Emu(60000), COLORS["أزرق"])
    info = f"{common['الجهة']}  |  {common['الحدث']}  |  {common['التاريخ']}"
    add_textbox(slide, info, Emu(300000), Emu(4000000),
                Emu(8544000), Emu(500000), size=13,
                color=COLORS["أزرق_فاتح"], align=PP_ALIGN.CENTER)

    # تحذير سري
    if common["سري"]:
        add_textbox(slide, "🔒 سري - للاستخدام الرسمي فقط",
                    Emu(300000), Emu(4600000), Emu(8544000), Emu(300000),
                    size=12, color=COLORS["أحمر"], align=PP_ALIGN.CENTER)

def build_bullets_slide(prs, slide_data, n, total):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_bg(slide, COLORS["خلفية"])

    # شريط العنوان
    add_rect(slide, Emu(0), Emu(0), Emu(9144000), Emu(900000), COLORS["أزرق_داكن"])
    icon  = slide_data.get("الأيقونة", "")
    title = f"{icon}  {slide_data['العنوان']}" if icon else slide_data["العنوان"]
    add_textbox(slide, title, Emu(200000), Emu(80000),
                Emu(8744000), Emu(730000),
                size=24, bold=True, color=COLORS["أبيض"])

    # خط فاصل
    add_rect(slide, Emu(200000), Emu(980000),
             Emu(8744000), Emu(50000), COLORS["أزرق"])

    # النقاط
    points = slide_data["النقاط"]
    y_start = Emu(1100000)
    spacing = Emu(530000) if len(points) <= 5 else Emu(440000)

    for i, point in enumerate(points):
        y = y_start + i * spacing
        # دائرة الرقم
        add_rect(slide, Emu(7900000), y + Emu(60000),
                 Emu(360000), Emu(360000), COLORS["أزرق"])
        add_textbox(slide, str(i+1),
                    Emu(7900000), y + Emu(60000),
                    Emu(360000), Emu(360000),
                    size=13, bold=True, color=COLORS["أبيض"],
                    align=PP_ALIGN.CENTER)
        # النص
        add_textbox(slide, f"  {point}",
                    Emu(200000), y, Emu(7600000), Emu(480000),
                    size=16, color=COLORS["نص"])

    add_slide_number(slide, n, total)

def build_colored_bullets_slide(prs, slide_data, n, total):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_bg(slide, COLORS["خلفية"])

    # شريط العنوان
    add_rect(slide, Emu(0), Emu(0), Emu(9144000), Emu(900000), COLORS["أزرق_داكن"])
    icon  = slide_data.get("الأيقونة", "")
    title = f"{icon}  {slide_data['العنوان']}" if icon else slide_data["العنوان"]
    add_textbox(slide, title, Emu(200000), Emu(80000),
                Emu(8744000), Emu(730000),
                size=24, bold=True, color=COLORS["أبيض"])
    add_rect(slide, Emu(200000), Emu(980000),
             Emu(8744000), Emu(50000), COLORS["أزرق"])

    points  = slide_data["النقاط"]
    y_start = Emu(1100000)
    spacing = Emu(540000) if len(points) <= 4 else Emu(460000)

    for i, (text, color_key) in enumerate(points):
        y    = y_start + i * spacing
        fg, bg = BULLET_COLORS.get(color_key, (COLORS["أزرق"], COLORS["أزرق_فاتح"]))
        # خلفية النقطة
        add_rect(slide, Emu(200000), y, Emu(8744000), Emu(480000), bg)
        # شريط اللون
        add_rect(slide, Emu(8700000), y, Emu(244000), Emu(480000), fg)
        # النص
        add_textbox(slide, f"  {text}",
                    Emu(200000), y, Emu(8480000), Emu(480000),
                    size=16, color=COLORS["نص"])

    add_slide_number(slide, n, total)

def build_numbered_slide(prs, slide_data, n, total):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_bg(slide, COLORS["خلفية"])

    add_rect(slide, Emu(0), Emu(0), Emu(9144000), Emu(900000), COLORS["أزرق_داكن"])
    icon  = slide_data.get("الأيقونة", "")
    title = f"{icon}  {slide_data['العنوان']}" if icon else slide_data["العنوان"]
    add_textbox(slide, title, Emu(200000), Emu(80000),
                Emu(8744000), Emu(730000),
                size=24, bold=True, color=COLORS["أبيض"])
    add_rect(slide, Emu(200000), Emu(980000),
             Emu(8744000), Emu(50000), COLORS["أزرق"])

    points  = slide_data["النقاط"]
    y_start = Emu(1120000)
    spacing = Emu(560000) if len(points) <= 4 else Emu(460000)

    for i, point in enumerate(points):
        y  = y_start + i * spacing
        bg = COLORS["أزرق_فاتح"] if i % 2 == 0 else COLORS["أبيض"]
        add_rect(slide, Emu(200000), y, Emu(8744000), Emu(490000), bg)
        # رقم دائري
        add_rect(slide, Emu(7980000), y + Emu(65000),
                 Emu(380000), Emu(380000), COLORS["أزرق_داكن"])
        add_textbox(slide, str(i+1),
                    Emu(7980000), y + Emu(65000),
                    Emu(380000), Emu(380000),
                    size=14, bold=True, color=COLORS["أبيض"],
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, f"  {point}",
                    Emu(200000), y, Emu(7680000), Emu(490000),
                    size=15, color=COLORS["نص"])

    add_slide_number(slide, n, total)

def build_cards_slide(prs, slide_data, n, total):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_bg(slide, COLORS["خلفية"])

    add_rect(slide, Emu(0), Emu(0), Emu(9144000), Emu(900000), COLORS["أزرق_داكن"])
    icon  = slide_data.get("الأيقونة", "")
    title = f"{icon}  {slide_data['العنوان']}" if icon else slide_data["العنوان"]
    add_textbox(slide, title, Emu(200000), Emu(80000),
                Emu(8744000), Emu(730000),
                size=24, bold=True, color=COLORS["أبيض"])
    add_rect(slide, Emu(200000), Emu(980000),
             Emu(8744000), Emu(50000), COLORS["أزرق"])

    cards   = slide_data["البطاقات"]
    n_cards = len(cards)
    card_w  = Emu(8744000 // n_cards - 100000)
    gap     = Emu(150000)
    x_start = Emu(200000)

    for i, card in enumerate(cards):
        x  = x_start + i * (card_w + gap)
        fg, bg = CARD_COLORS.get(card["اللون"], (COLORS["أزرق"], COLORS["أزرق_فاتح"]))

        # بطاقة رئيسية
        add_rect(slide, x, Emu(1150000), card_w, Emu(3300000), bg,
                 line_color=fg)
        # رأس البطاقة
        add_rect(slide, x, Emu(1150000), card_w, Emu(600000), fg)
        add_textbox(slide, card["العنوان"],
                    x, Emu(1150000), card_w, Emu(600000),
                    size=16, bold=True, color=COLORS["أبيض"],
                    align=PP_ALIGN.CENTER)
        # القيمة
        add_textbox(slide, card["القيمة"],
                    x, Emu(2000000), card_w, Emu(800000),
                    size=28, bold=True, color=fg,
                    align=PP_ALIGN.CENTER)
        # الوصف
        add_textbox(slide, card["الوصف"],
                    x, Emu(3000000), card_w, Emu(400000),
                    size=13, color=COLORS["رمادي"],
                    align=PP_ALIGN.CENTER)

    add_slide_number(slide, n, total)

def build_closing_slide(prs, common):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_bg(slide, COLORS["أزرق_داكن"])

    add_rect(slide, Emu(0), Emu(4300000), Emu(9144000), Emu(200000), COLORS["أزرق"])
    add_textbox(slide, "شكراً لحسن الاستماع",
                Emu(300000), Emu(1500000), Emu(8544000), Emu(1000000),
                size=36, bold=True, color=COLORS["أبيض"],
                align=PP_ALIGN.CENTER)
    add_textbox(slide, common["الجهة"],
                Emu(300000), Emu(2700000), Emu(8544000), Emu(600000),
                size=18, color=COLORS["أزرق_فاتح"],
                align=PP_ALIGN.CENTER)
    add_textbox(slide, f"{common['الحدث']}  |  {common['التاريخ']}",
                Emu(300000), Emu(3400000), Emu(8544000), Emu(500000),
                size=14, color=COLORS["رمادي"],
                align=PP_ALIGN.CENTER)

# ========== المولّد الرئيسي ==========
def generate(pres_type=PRESENTATION_TYPE):
    pres_data = PRESENTATIONS.get(pres_type)
    if not pres_data:
        print(f"نوع العرض غير موجود: {pres_type}")
        return None

    prs = Presentation()
    prs.slide_width  = Emu(9144000)
    prs.slide_height = Emu(5143500)

    slides     = pres_data["الشرائح"]
    total      = len(slides) + 2   # غلاف + شرائح + ختام

    # شريحة الغلاف
    build_title_slide(prs, pres_data, COMMON)

    # شرائح المحتوى
    for i, slide_data in enumerate(slides, 1):
        ntype = slide_data["النوع"]
        if ntype == "نقاط":
            build_bullets_slide(prs, slide_data, i, total-2)
        elif ntype == "نقاط_ملونة":
            build_colored_bullets_slide(prs, slide_data, i, total-2)
        elif ntype == "مرقّم":
            build_numbered_slide(prs, slide_data, i, total-2)
        elif ntype == "بطاقات":
            build_cards_slide(prs, slide_data, i, total-2)

    # شريحة الختام
    build_closing_slide(prs, COMMON)

    return prs

# ========== التشغيل ==========
if __name__ == "__main__":
    print(f"جارٍ توليد عرض: {PRESENTATION_TYPE}...")

    prs = generate(PRESENTATION_TYPE)
    if prs:
        name    = PRESENTATIONS[PRESENTATION_TYPE]["العنوان"][:20].strip()
        outfile = f"tools/pptx_generator/output_PPTX_{PRESENTATION_TYPE}.pptx"
        prs.save(outfile)
        print(f"✓ تم إنشاء العرض: {outfile}")
        print(f"  عدد الشرائح: {len(prs.slides)}")
        print(f"\nلتغيير العرض افتح pptx_data.py وعدّل PRESENTATION_TYPE إلى:")
        print("  'مادة_سادسة' | 'جرد'")
