"""
النموذج الثالث: ppt-generation (keithmcnulty)
استبدال بيانات رقمية في عرض PowerPoint موجود برمجياً
الفائدة: تحديث أرقام عرض الجرد السنوي دون لمس التصميم
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import copy, os

# ========== تعريف خرائط الاستبدال ==========
# القالب القديم -> البيانات الجديدة
# أضف هنا أي متغير تريد تحديثه في نسخة الجرد القادمة

replacements = {
    # أرقام الجرد
    "{{إجمالي_الانبعاثات}}": "١٨٥.٤",
    "{{انبعاثات_الطاقة}}":   "١٤٢.٣",
    "{{انبعاثات_الصناعة}}":  "٢٨.١",
    "{{انبعاثات_النفايات}}":  "١٥.٠",
    "{{السنة}}":              "٢٠٢٤",
    "{{نسبة_التغيير}}":       "+٣.٢٪",
    # معلومات التقرير
    "{{اسم_المعد}}":          "قسم سياسات المناخ",
    "{{تاريخ_الإصدار}}":      "مايو ٢٠٢٥",
}

def replace_text_in_shape(shape, replacements):
    """استبدال النصوص في شكل مع الحفاظ على التنسيق"""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            for old, new in replacements.items():
                if old in run.text:
                    run.text = run.text.replace(old, new)

def process_presentation(input_path, output_path, replacements):
    """معالجة الملف واستبدال جميع المتغيرات"""
    prs = Presentation(input_path)
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for old, new in replacements.items():
                            if old in run.text:
                                run.text = run.text.replace(old, new)
                                count += 1
            # معالجة الجداول
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                for old, new in replacements.items():
                                    if old in run.text:
                                        run.text = run.text.replace(old, new)
                                        count += 1
    prs.save(output_path)
    return count

# ========== إنشاء ملف قالب تجريبي أولاً ==========
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Emu(9144000)
prs.slide_height = Emu(5143500)
blank = prs.slide_layouts[6]

# شريحة القالب بمتغيرات مدمجة
slide = prs.slides.add_slide(blank)
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFF)

def add_shape_text(slide, text, l, t, w, h, size=14, bold=False, color=None):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para._p.get_or_add_pPr().set('rtl', '1')
    para.alignment = PP_ALIGN.RIGHT
    run = para.add_run()
    run.text = text
    run.font.name = "Simplified Arabic"
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color

# عنوان
add_shape_text(slide, "تقرير الجرد الوطني - سنة {{السنة}}",
               Emu(200000), Emu(200000), Emu(8744000), Emu(700000),
               size=26, bold=True, color=RGBColor(0x1F, 0x49, 0x7D))

# صف البيانات
add_shape_text(slide, "إجمالي الانبعاثات: {{إجمالي_الانبعاثات}} مليون طن",
               Emu(200000), Emu(1100000), Emu(8744000), Emu(500000),
               size=18, color=RGBColor(0x26, 0x26, 0x26))

add_shape_text(slide, "قطاع الطاقة: {{انبعاثات_الطاقة}} | الصناعة: {{انبعاثات_الصناعة}} | النفايات: {{انبعاثات_النفايات}}",
               Emu(200000), Emu(1700000), Emu(8744000), Emu(500000),
               size=16, color=RGBColor(0x40, 0x40, 0x40))

add_shape_text(slide, "نسبة التغيير عن العام السابق: {{نسبة_التغيير}}",
               Emu(200000), Emu(2400000), Emu(8744000), Emu(500000),
               size=16, color=RGBColor(0x1F, 0x49, 0x7D))

add_shape_text(slide, "معد التقرير: {{اسم_المعد}}    |    تاريخ الإصدار: {{تاريخ_الإصدار}}",
               Emu(200000), Emu(4500000), Emu(8744000), Emu(400000),
               size=12, color=RGBColor(0x70, 0x70, 0x70))

# حفظ القالب
template_path = "/home/claude/skills-repo/demos/template_قالب_الجرد.pptx"
prs.save(template_path)
print(f"تم إنشاء القالب: {template_path}")

# ========== تطبيق الاستبدال ==========
output_path = "/home/claude/skills-repo/demos/output_demo3_جرد_محدّث.pptx"
count = process_presentation(template_path, output_path, replacements)
print(f"تم الاستبدال في {count} موضع")
print(f"تم إنشاء الملف: {output_path}")
