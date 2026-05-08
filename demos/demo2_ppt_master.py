"""
النموذج الثاني: ppt-master / python-pptx
توليد عرض PowerPoint عربي قابل للتعديل الكامل
الفائدة: إنتاج عروض المفاوضات والمؤتمرات من بيانات نصية
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ========== بيانات العرض ==========
presentation_data = {
    "العنوان": "آليات المادة السادسة من اتفاقية باريس",
    "العنوان_الفرعي": "فرص وتحديات في ضوء مخرجات COP29",
    "الجهة": "قسم سياسات المناخ",
    "الشرائح": [
        {
            "العنوان": "ما هي المادة السادسة؟",
            "النقاط": [
                "تتيح التعاون الطوعي بين الدول لتحقيق أهداف المناخ",
                "تشمل آليتين رئيسيتين: 6.2 و6.4",
                "6.2: نتائج التخفيف المحوّلة دولياً (ITMOs)",
                "6.4: آلية الأمم المتحدة للكربون الدولية",
            ]
        },
        {
            "العنوان": "أبرز مخرجات COP29",
            "النقاط": [
                "اعتماد معايير التحقق للمادة 6.4",
                "إطلاق سجل المعاملات الدولي",
                "اتفاق 30 دولة على صفقات ITMOs أولية",
                "تحديات لا تزال قائمة في تجنب الازدواجية",
            ]
        },
        {
            "العنوان": "موقف مجموعة الدول العربية",
            "النقاط": [
                "دعم آليات السوق مع ضمان التنمية المستدامة",
                "التأكيد على حق الدول النامية في التنمية",
                "مطالبة بحصة من العائدات لصندوق التكيف",
                "رفض الشروط المسبقة على استخدام الكربون الأحفوري",
            ]
        },
    ]
}

# ========== ألوان القالب ==========
COLOR_DARK    = RGBColor(0x1F, 0x49, 0x7D)
COLOR_ACCENT  = RGBColor(0x2E, 0x75, 0xB6)
COLOR_LIGHT   = RGBColor(0xDE, 0xEB, 0xF7)
COLOR_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TEXT    = RGBColor(0x26, 0x26, 0x26)
FONT          = "Simplified Arabic"

prs = Presentation()
prs.slide_width  = Emu(9144000)   # 10 inches
prs.slide_height = Emu(5143500)   # 7.5 inches

# ========== دالة مساعدة: إعداد RTL ==========
def make_rtl(tf):
    for para in tf.paragraphs:
        pPr = para._p.get_or_add_pPr()
        algn = etree.SubElement(pPr, qn('a:buNone')) if False else None
        # RTL
        pPr.set('rtl', '1')
        para.alignment = PP_ALIGN.RIGHT

def add_text_rtl(tf, text, size=18, bold=False, color=None):
    para = tf.add_paragraph()
    pPr = para._p.get_or_add_pPr()
    pPr.set('rtl', '1')
    para.alignment = PP_ALIGN.RIGHT
    run = para.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return para

def add_bullet_rtl(tf, text, size=16, color=None):
    para = tf.add_paragraph()
    pPr = para._p.get_or_add_pPr()
    pPr.set('rtl', '1')
    para.alignment = PP_ALIGN.RIGHT
    para.level = 0
    run = para.add_run()
    run.text = "◀  " + text
    run.font.name = FONT
    run.font.size = Pt(size)
    if color:
        run.font.color.rgb = color
    return para

def add_textbox(slide, text, left, top, width, height, size=18, bold=False,
                color=COLOR_TEXT, bg=None, align=PP_ALIGN.RIGHT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    if bg:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg
    para = tf.paragraphs[0]
    pPr = para._p.get_or_add_pPr()
    pPr.set('rtl', '1')
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

# ========== شريحة العنوان ==========
blank = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(blank)
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = COLOR_DARK

# شريط زخرفي
shape = slide.shapes.add_shape(1, Emu(0), Emu(4400000), Emu(9144000), Emu(300000))
shape.fill.solid()
shape.fill.fore_color.rgb = COLOR_ACCENT
shape.line.fill.background()

# العنوان
add_textbox(slide, presentation_data["العنوان"],
            Emu(300000), Emu(1500000), Emu(8544000), Emu(1200000),
            size=32, bold=True, color=COLOR_WHITE)

# العنوان الفرعي
add_textbox(slide, presentation_data["العنوان_الفرعي"],
            Emu(300000), Emu(2800000), Emu(8544000), Emu(800000),
            size=20, color=COLOR_LIGHT)

# اسم الجهة
add_textbox(slide, presentation_data["الجهة"],
            Emu(300000), Emu(4500000), Emu(8544000), Emu(500000),
            size=14, color=COLOR_WHITE)

# ========== شرائح المحتوى ==========
for i, slide_data in enumerate(presentation_data["الشرائح"], 1):
    slide = prs.slides.add_slide(blank)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFF)

    # شريط العنوان العلوي
    header = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(9144000), Emu(950000))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_DARK
    header.line.fill.background()

    # عنوان الشريحة
    add_textbox(slide, slide_data["العنوان"],
                Emu(200000), Emu(80000), Emu(8744000), Emu(780000),
                size=24, bold=True, color=COLOR_WHITE)

    # خط فاصل زرقاء فاتحة
    sep = slide.shapes.add_shape(1, Emu(200000), Emu(1050000), Emu(8744000), Emu(60000))
    sep.fill.solid()
    sep.fill.fore_color.rgb = COLOR_ACCENT
    sep.line.fill.background()

    # منطقة النقاط
    txBox = slide.shapes.add_textbox(
        Emu(200000), Emu(1200000), Emu(8744000), Emu(3600000))
    tf = txBox.text_frame
    tf.word_wrap = True

    for j, bullet in enumerate(slide_data["النقاط"]):
        if j == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        pPr = para._p.get_or_add_pPr()
        pPr.set('rtl', '1')
        para.alignment = PP_ALIGN.RIGHT
        para.space_before = Pt(8)
        run = para.add_run()
        run.text = "◀  " + bullet
        run.font.name = FONT
        run.font.size = Pt(18)
        run.font.color.rgb = COLOR_TEXT

    # رقم الشريحة
    add_textbox(slide, f"{i} / {len(presentation_data['الشرائح'])}",
                Emu(100000), Emu(4700000), Emu(500000), Emu(300000),
                size=11, color=COLOR_ACCENT)

# ========== حفظ الملف ==========
output_path = "/home/claude/skills-repo/demos/output_demo2_عرض_المادة_السادسة.pptx"
prs.save(output_path)
print(f"تم إنشاء الملف: {output_path}")
